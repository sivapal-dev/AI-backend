import os
import re
import json
import subprocess
import tempfile
from typing import Optional, Dict, Any
from pathlib import Path
from datetime import datetime, timezone
from openai import OpenAI
from config import get_settings
from services.file_storage import generate_file_name, get_project_dir, get_file_size
import logging

logger = logging.getLogger(__name__)
settings = get_settings()

# ─── Backend directory (where node_modules lives) ─────────────────────────────
BACKEND_DIR = Path(__file__).parent.parent


import html

def clean_text_content(text: str) -> str:
    if not isinstance(text, str):
        return str(text) if text is not None else ""
    # Decode HTML entities
    text = html.unescape(text)
    # Strip HTML tags
    text = re.sub(r'<[^>]+>', '', text)
    # Strip leading markdown heading symbols from normal texts
    text = re.sub(r'^#+\s+', '', text)
    # Strip leading list symbols from paragraph/heading lines
    text = re.sub(r'^[-\*\+]\s+', '', text)
    # Strip backticks
    text = text.replace("`", "")
    return text.strip()


def clean_json_data(data: Any) -> Any:
    if isinstance(data, dict):
        cleaned = {}
        for k, v in data.items():
            if k in ["title", "subtitle", "author", "date", "text"]:
                cleaned[k] = clean_text_content(v)
            elif k == "items" and isinstance(v, list):
                cleaned[k] = [clean_text_content(item) for item in v]
            elif k == "headers" and isinstance(v, list):
                cleaned[k] = [clean_text_content(item) for item in v]
            elif k == "rows" and isinstance(v, list):
                cleaned[k] = [[clean_text_content(cell) for cell in row] if isinstance(row, list) else [] for row in v]
            else:
                cleaned[k] = clean_json_data(v)
        return cleaned
    elif isinstance(data, list):
        return [clean_json_data(item) for item in data]
    return data


def parse_malformed_ai_response(raw_text: str, prompt: str) -> dict:
    lines = []
    for line in raw_text.splitlines():
        line = line.strip()
        if not line or line in ["{", "}", "[", "]", "],", "},"]:
            continue
        match = re.match(r'^"[a-zA-Z0-9_\-]+"\s*:\s*"(.*)"\s*,?$', line)
        if match:
            line = match.group(1)
        elif line.startswith('"') and line.endswith('"'):
            line = line[1:-1]
        elif line.startswith('"') and line.endswith('",'):
            line = line[1:-2]
        line = line.strip()
        if line:
            lines.append(line)
            
    sections = []
    current_section = {"title": "Main Content", "content": []}
    current_list_items = []
    current_list_style = "bullet"
    
    def flush_list():
        nonlocal current_list_items
        if current_list_items:
            current_section["content"].append({
                "type": "list",
                "style": current_list_style,
                "items": current_list_items
            })
            current_list_items = []

    for line in lines:
        cleaned = clean_text_content(line)
        if not cleaned:
            continue
            
        if line.startswith("#") or (line.isupper() and len(cleaned) < 60):
            flush_list()
            if current_section["content"]:
                sections.append(current_section)
                current_section = {"title": cleaned, "content": []}
            else:
                current_section["title"] = cleaned
        elif re.match(r'^\d+\.\s+', line):
            if current_list_items and current_list_style != "number":
                flush_list()
            current_list_style = "number"
            item_text = re.sub(r'^\d+\.\s+', '', line)
            current_list_items.append(clean_text_content(item_text))
        elif line.startswith("- ") or line.startswith("* ") or line.startswith("+ "):
            if current_list_items and current_list_style != "bullet":
                flush_list()
            current_list_style = "bullet"
            item_text = re.sub(r'^[-\*\+]\s+', '', line)
            current_list_items.append(clean_text_content(item_text))
        else:
            flush_list()
            current_section["content"].append({
                "type": "paragraph",
                "text": cleaned
            })
            
    flush_list()
    if current_section["content"] or len(sections) == 0:
        sections.append(current_section)
        
    return {
        "title": clean_text_content(prompt[:60]),
        "subtitle": "AI Generated Document",
        "author": "By8flow AI",
        "sections": sections
    }


def parse_malformed_ai_response_for_ppt(raw_text: str, prompt: str) -> dict:
    lines = []
    for line in raw_text.splitlines():
        line = line.strip()
        if not line or line in ["{", "}", "[", "]", "],", "},"]:
            continue
        match = re.match(r'^"[a-zA-Z0-9_\-]+"\s*:\s*"(.*)"\s*,?$', line)
        if match:
            line = match.group(1)
        elif line.startswith('"') and line.endswith('"'):
            line = line[1:-1]
        elif line.startswith('"') and line.endswith('",'):
            line = line[1:-2]
        line = line.strip()
        if line:
            lines.append(line)

    slides = []
    current_slide = {"title": "Overview", "layout": "content", "content": []}
    for line in lines:
        cleaned = clean_text_content(line)
        if not cleaned:
            continue
        if len(current_slide["content"]) >= 5:
            slides.append(current_slide)
            current_slide = {"title": "Key Details", "layout": "content", "content": []}
        current_slide["content"].append(cleaned)
    if current_slide["content"] or len(slides) == 0:
        slides.append(current_slide)

    return {
        "title": clean_text_content(prompt[:60]),
        "subtitle": "AI Generated Presentation",
        "author": "By8flow AI",
        "theme": "corporate",
        "slides": slides
    }


def _build_pptx_prompt(prompt: str, project_context: str) -> str:
    schema = """{
  "title": "Presentation Title",
  "subtitle": "Presentation Subtitle",
  "author": "Author Name",
  "date": "Date of Presentation (optional)",
  "theme": "corporate" | "modern" | "creative",
  "slides": [
    {
      "title": "Slide Title",
      "layout": "title" | "content" | "two_columns" | "table" | "closing",
      "content": [
        "Bullet item 1",
        "Bullet item 2"
      ],
      "content_left": [
        "Left column bullet (for two_columns layout)"
      ],
      "content_right": [
        "Right column bullet (for two_columns layout)"
      ],
      "table": {
        "headers": ["Col 1", "Col 2"],
        "rows": [
          ["Row 1 Col 1", "Row 1 Col 2"],
          ["Row 2 Col 1", "Row 2 Col 2"]
        ]
      }
    }
  ]
}"""
    return f"""You are a world-class PowerPoint presentation designer. Generate a COMPLETE presentation structure strictly adhering to the JSON schema below.

DESIGN PHILOSOPHY:
- Minimal, confident, data-rich
- Choose a theme from: "corporate", "modern", "creative"

SLIDE STRUCTURE (minimum 6 slides, maximum 12):
1. Title slide (layout: "title")
2. Executive summary / Agenda (layout: "content")
3. Content slides (2-5 slides) (layout: "content" or "two_columns")
4. Table / data slide (layout: "table")
5. Conclusion / Key Takeaways (layout: "content")
6. Closing slide (layout: "closing")

ABSOLUTE CRITICAL RULES:
1. Output ONLY the raw JSON object. Do NOT wrap it in markdown fences (no ```json).
2. The output must be valid, well-formed JSON. Do not include comments, explanations, or prose.

JSON Schema format:
{schema}

USER REQUEST:
{prompt.strip()[:3000]}

PROJECT CONTEXT:
{project_context[:1000] if project_context else "None"}

Generate ONLY the raw JSON content.
"""


def _build_docx_prompt(prompt: str, project_context: str) -> str:
    schema = """{
  "title": "Document Title",
  "subtitle": "Document Subtitle",
  "author": "Author Name",
  "date": "Date of Document (optional)",
  "sections": [
    {
      "title": "Section Title",
      "content": [
        { "type": "paragraph", "text": "Prose text here. You can use standard markdown inline **bold** or *italics* formatting." },
        { "type": "heading", "level": 1 | 2 | 3, "text": "Heading text" },
        { "type": "list", "style": "bullet" | "number", "items": ["Item 1", "Item 2"] },
        { "type": "code", "text": "Code snippet..." },
        {
          "type": "table",
          "headers": ["Col 1", "Col 2"],
          "rows": [
            ["Val 1", "Val 2"],
            ["Val 3", "Val 4"]
          ]
        }
      ]
    }
  ]
}"""
    return f"""You are a world-class document designer. Generate a COMPLETE document structure strictly adhering to the JSON schema below.

DOCUMENT STRUCTURE (minimum):
1. Cover page metadata (title, subtitle, author)
2. Table of Contents (as formatted lists/headings)
3. Executive Summary
4. Body sections with paragraphs, headings, bullet lists, and tables
5. Conclusion / Recommendations

ABSOLUTE CRITICAL RULES:
1. Output ONLY the raw JSON object. Do NOT wrap it in markdown fences (no ```json).
2. The output must be valid, well-formed JSON. Do not include comments, explanations, or prose.

JSON Schema format:
{schema}

USER REQUEST:
{prompt.strip()[:3000]}

PROJECT CONTEXT:
{project_context[:1000] if project_context else "None"}

Generate ONLY the raw JSON content.
"""


def _strip_and_sanitize_json(raw_json: str) -> str:
    """
    Clean AI-generated JSON output:
    1. Strip markdown code fences (```json or ``` blocks)
    2. Remove any leading/trailing prose
    """
    json_str = raw_json.strip()

    # Remove markdown fences — handle ```json, ``` etc.
    json_str = re.sub(r'^```(?:json)?\s*', '', json_str, flags=re.MULTILINE)
    json_str = re.sub(r'```\s*$', '', json_str, flags=re.MULTILINE)
    json_str = json_str.strip()

    # If there's prose before the first '{', strip it
    first_brace = json_str.find('{')
    if first_brace != -1:
        json_str = json_str[first_brace:]

    # If there's prose after the last '}', strip it
    last_brace = json_str.rfind('}')
    if last_brace != -1:
        json_str = json_str[:last_brace + 1]

    return json_str.strip()


async def generate_document(
    prompt: str,
    format_type: str,
    project_id: Optional[str] = None,
    project_context: str = "",
) -> Dict[str, Any]:
    """
    Generate a document (DOC or PPT) from a prompt using AI.
    Returns dict with file_path, file_name, pdf_path (for PPT), size_bytes.
    """
    format_ext = "docx" if format_type == "doc" else "pptx"
    file_name = generate_file_name(None, format_ext)
    project_dir = get_project_dir(project_id)
    output_path = str(project_dir / file_name)

    # Build the AI prompt
    if format_type == "ppt":
        ai_prompt = _build_pptx_prompt(prompt, project_context)
    else:
        ai_prompt = _build_docx_prompt(prompt, project_context)

    # Call OpenRouter AI
    try:
        client = OpenAI(
            base_url="https://openrouter.ai/api/v1",
            api_key=settings.openrouter_api_key,
        )
        response = client.chat.completions.create(
            model=settings.openrouter_default_model,
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You are a structured data generator that outputs ONLY raw, valid JSON. "
                        "Never output markdown fences, explanations, comments, or prose. "
                        "Output ONLY the valid JSON object, nothing else."
                    )
                },
                {"role": "user", "content": ai_prompt},
            ],
            temperature=0.5,  # Lower temperature for more stable JSON formatting
            max_tokens=14000,
        )
        raw_output = response.choices[0].message.content or ""
        json_content = _strip_and_sanitize_json(raw_output)
    except Exception as e:
        logger.error(f"AI generation failed: {e}")
        raise Exception(f"AI generation failed: {str(e)}")

    if not json_content:
        raise Exception("AI returned empty content. Please try again with a different prompt.")

    # Validate JSON parsing in Python first
    try:
        parsed_data = json.loads(json_content, strict=False)
        # Ensure we have at least standard structure
        if format_type == "doc":
            if "title" not in parsed_data:
                parsed_data["title"] = prompt[:50]
            if "sections" not in parsed_data or not parsed_data["sections"]:
                parsed_data["sections"] = [{"title": "Content", "content": [{"type": "paragraph", "text": prompt}]}]
            parsed_data = clean_json_data(parsed_data)
        else:
            if "title" not in parsed_data:
                parsed_data["title"] = prompt[:50]
            if "slides" not in parsed_data or not parsed_data["slides"]:
                parsed_data["slides"] = [{"title": "Content", "layout": "content", "content": [prompt]}]
            parsed_data = clean_json_data(parsed_data)
        # Re-serialize clean JSON
        json_content = json.dumps(parsed_data)
    except Exception as e:
        logger.error(f"AI returned malformed JSON: {e}. Raw content: {raw_output[:500]}")
        # Try to fallback to a basic document if JSON parsing fails completely
        if format_type == "doc":
            fallback_data = parse_malformed_ai_response(raw_output, prompt)
        else:
            fallback_data = parse_malformed_ai_response_for_ppt(raw_output, prompt)
        json_content = json.dumps(fallback_data)

    # Write JSON data to temp file for the builder script
    temp_dir = tempfile.mkdtemp()
    temp_json = os.path.join(temp_dir, "data.json")
    with open(temp_json, "w", encoding="utf-8") as f:
        f.write(json_content)

    builder_path = str(BACKEND_DIR / "services" / "document_builder.js")
    logger.info(f"Executing document builder script: {builder_path}")

    # Build node environment
    node_env = os.environ.copy()
    node_modules_path = str(BACKEND_DIR / "node_modules")
    node_env["NODE_PATH"] = node_modules_path

    # Run Node document builder CLI
    try:
        result = subprocess.run(
            ["node", builder_path, temp_json, output_path, format_type],
            capture_output=True,
            text=True,
            timeout=90,
            cwd=str(BACKEND_DIR),
            env=node_env,
        )
        if result.stdout:
            logger.info(f"Builder stdout: {result.stdout.strip()[:500]}")
        if result.stderr:
            logger.warning(f"Builder stderr: {result.stderr.strip()[:500]}")

        if result.returncode != 0:
            err_msg = result.stderr.strip()[:800] if result.stderr else "Unknown error"
            logger.error(f"Builder execution failed (returncode={result.returncode}): {err_msg}")
            raise Exception(f"Document compilation failed: {err_msg}")

    except subprocess.TimeoutExpired:
        if os.path.exists(output_path):
            try:
                os.remove(output_path)
            except Exception:
                pass
        raise Exception("Document compilation timed out (90s limit). Try a shorter prompt.")
    except FileNotFoundError:
        raise Exception("Node.js not found. Please ensure Node.js is installed.")
    finally:
        # Cleanup temp file
        try:
            import shutil
            shutil.rmtree(temp_dir, ignore_errors=True)
        except Exception:
            pass

    # Verify file was created
    if not os.path.exists(output_path):
        raise Exception(
            f"Document file was not created at expected path: {output_path}. "
            f"Please try again."
        )

    size_bytes = get_file_size(output_path)
    logger.info(f"Document generated successfully: {output_path} ({size_bytes} bytes)")

    # For PPT: convert to PDF for preview
    pdf_path = None
    if format_type == "ppt":
        pdf_path = await _convert_pptx_to_pdf(output_path, project_dir)

    return {
        "file_path": output_path,
        "file_name": file_name,
        "pdf_path": pdf_path,
        "size_bytes": size_bytes,
    }


async def _convert_pptx_to_pdf(pptx_path: str, output_dir: Path) -> Optional[str]:
    """Convert PPTX to PDF using LibreOffice (headless) or python-pptx fallback."""
    pdf_name = Path(pptx_path).stem + ".pdf"
    pdf_path = str(output_dir / pdf_name)

    # Try LibreOffice first (best quality)
    try:
        result = subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf", "--outdir", str(output_dir), pptx_path],
            capture_output=True,
            text=True,
            timeout=120,
        )
        if result.returncode == 0 and os.path.exists(pdf_path):
            logger.info(f"LibreOffice converted PPTX to PDF: {pdf_path}")
            return pdf_path
        else:
            # Clean up partial/corrupted files on failure (W200)
            if os.path.exists(pdf_path):
                try:
                    os.remove(pdf_path)
                except Exception:
                    pass
    except subprocess.TimeoutExpired as e:
        logger.warning(f"LibreOffice conversion timed out: {e}")
        if os.path.exists(pdf_path):
            try:
                os.remove(pdf_path)
            except Exception:
                pass
    except Exception as e:
        logger.warning(f"LibreOffice not available or failed: {e}")
        if os.path.exists(pdf_path):
            try:
                os.remove(pdf_path)
            except Exception:
                pass

    # Fallback: try python-pptx + reportlab to create a simple PDF preview
    try:
        from pptx import Presentation
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas

        prs = Presentation(pptx_path)
        c = canvas.Canvas(pdf_path, pagesize=letter)
        width, height = letter

        # Get slide dimensions to scale positions (W201)
        slide_w = prs.slide_width
        slide_h = prs.slide_height

        scale_x = width / slide_w if slide_w else 1.0
        scale_y = height / slide_h if slide_h else 1.0

        for i, slide in enumerate(prs.slides):
            if i > 0:
                c.showPage()
            
            # Draw slide background/header
            c.setFont("Helvetica-Bold", 8)
            c.setFillColorRGB(0.5, 0.5, 0.5)
            c.drawString(30, height - 20, f"By8flow Presentation Slide Preview — Slide {i + 1}")
            c.setStrokeColorRGB(0.8, 0.8, 0.8)
            c.line(30, height - 25, width - 30, height - 25)
            c.setFillColorRGB(0, 0, 0)  # reset fill

            # Parse elements with positions and sizes (W201)
            for shape in slide.shapes:
                if shape.has_text_frame:
                    # Calculate scaled coordinates
                    x = shape.left * scale_x if hasattr(shape, 'left') else 50
                    y = height - (shape.top * scale_y) - 20 if hasattr(shape, 'top') else (height - 50)
                    
                    # Ensure coordinates are within page boundaries
                    x = max(30, min(width - 100, x))
                    y = max(30, min(height - 50, y))

                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            # Determine font size from element attributes (W201)
                            font_size = 12
                            if para.font and para.font.size:
                                font_size = int(para.font.size.pt) if hasattr(para.font.size, 'pt') else 12
                                font_size = max(8, min(36, font_size))
                            
                            c.setFont("Helvetica", font_size)
                            # Wrap text if too long
                            chunk_size = 80
                            lines = [text[idx:idx+chunk_size] for idx in range(0, len(text), chunk_size)]
                            
                            for line in lines:
                                c.drawString(x, y, line)
                                y -= (font_size + 6)
                                
                                # Page break trigger if text height exceeds page boundaries (W202)
                                if y < 40:
                                    c.showPage()
                                    y = height - 50
                                    c.setFont("Helvetica", font_size)

            c.setFont("Helvetica", 10)
            c.drawString(width - 100, 20, f"Page {i + 1}")

        c.save()
        if os.path.exists(pdf_path):
            logger.info(f"Python fallback created PDF preview: {pdf_path}")
            return pdf_path
    except ImportError:
        logger.warning("python-pptx or reportlab not installed. PDF preview unavailable.")
        if os.path.exists(pdf_path):
            try:
                os.remove(pdf_path)
            except Exception:
                pass
    except Exception as e:
        logger.warning(f"Python PDF conversion failed: {e}")
        if os.path.exists(pdf_path):
            try:
                os.remove(pdf_path)
            except Exception:
                pass

    logger.warning("PDF conversion not available.")
    return None
